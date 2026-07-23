# -*- coding: utf-8 -*-
"""
PPTX 生成器：根据全链路结构化产出，生成 15 页标准可提报创意方案。
使用 python-pptx，遵循《工作流》第八节的标准页面结构。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FONT = "Microsoft YaHei"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _rgb(hex_str: str, default="#C8102E") -> RGBColor:
    h = (hex_str or default).lstrip("#")
    if len(h) != 6:
        h = default.lstrip("#")
    return RGBColor.from_string(h)


def _txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, space_after=6):
    """runs: list of (text, size, bold, color_hex) 或 单个 tuple -> 多段落。"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(runs, tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if isinstance(para, tuple):
            para = [para]
        for (text, size, bold, color) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = _rgb(color)
            # 中文字体（东亚）
            _set_east_asian(r, FONT)
    return tb


def _set_east_asian(run, font_name):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font_name)


def _rect(slide, l, t, w, h, fill_hex, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill_hex)
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _top_bar(slide, title, accent):
    _rect(slide, 0, 0, 13.333, 1.15, accent)
    _txt(slide, 0.6, 0.18, 12.1, 0.8,
         [(title, 26, True, "#FFFFFF")], anchor=MSO_ANCHOR.MIDDLE)


def _footer(slide, project, page_no):
    name = project["meta"].get("project_name", "") or "商业视频创意方案"
    _txt(slide, 0.6, 7.05, 9, 0.35,
         [(f"{name}  ·  创意提报方案", 9, False, "#9AA0A6")])
    _txt(slide, 11.8, 7.05, 1.0, 0.35,
         [(str(page_no), 9, False, "#9AA0A6")], align=PP_ALIGN.RIGHT)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = 空白


def _bullet(slide, l, t, w, h, items, accent, size=14, gap=8):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        # 项目符号
        rb = p.add_run(); rb.text = "▪ "; rb.font.size = Pt(size); rb.font.bold = True
        rb.font.color.rgb = _rgb(accent); rb.font.name = FONT
        _set_east_asian(rb, FONT)
        if isinstance(it, tuple):
            head, body = it
            rh = p.add_run(); rh.text = head; rh.font.size = Pt(size); rh.font.bold = True
            rh.font.color.rgb = _rgb("#202124"); rh.font.name = FONT; _set_east_asian(rh, FONT)
            if body:
                rr = p.add_run(); rr.text = "  " + body; rr.font.size = Pt(size); rr.font.bold = False
                rr.font.color.rgb = _rgb("#3C4043"); rr.font.name = FONT; _set_east_asian(rr, FONT)
        else:
            rr = p.add_run(); rr.text = it; rr.font.size = Pt(size); rr.font.bold = False
            rr.font.color.rgb = _rgb("#3C4043"); rr.font.name = FONT; _set_east_asian(rr, FONT)
    return tb


# ---------------------------------------------------------------------------
def build(project: dict, output_path: str) -> str:
    accent = project["meta"].get("accent_color", "#C8102E")
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    # 封面 P1
    _cover(prs, project, accent)
    # 目录 P2
    _toc(prs, project, accent)
    # P3 项目理解
    _p3_understanding(prs, project, accent)
    # P4 核心策略
    _p4_strategy(prs, project, accent)
    # P5 创意概念
    _p5_concept(prs, project, accent)
    # P6 故事梗概
    _p6_story(prs, project, accent)
    # P7-P9 分镜精选
    _p7_9_storyboard(prs, project, accent)
    # P10 视觉风格
    _p10_style(prs, project, accent)
    # P11 美术置景
    _p11_art(prs, project, accent)
    # P12 演员与角色
    _p12_actors(prs, project, accent)
    # P13 制作排期
    _p13_schedule(prs, project, accent)
    # P14 团队与优势
    _p14_team(prs, project, accent)
    # P15 结尾
    _p15_end(prs, project, accent)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


# ---- P1 封面 ---------------------------------------------------------------
def _cover(prs, project, accent):
    s = _blank(prs)
    _rect(s, 0, 0, 13.333, 7.5, accent)
    _rect(s, 0, 5.3, 13.333, 0.06, "#FFFFFF")
    meta = project["meta"]
    dirs = project["stage2"].get("directions", [])
    idx = min(project["stage2"].get("chosen", 0), max(len(dirs) - 1, 0))
    dname = dirs[idx].get("name", "") if idx < len(dirs) else ""
    ptitle = project["stage8"].get("ppt_title") or meta.get("project_name") or "商业视频创意方案"
    psub = project["stage8"].get("ppt_subtitle") or (f"创意方向 · 「{dname}」" if dname else "创意提报方案")

    _txt(s, 0.9, 1.6, 11.5, 1.0, [(ptitle, 40, True, "#FFFFFF")])
    _txt(s, 0.9, 2.9, 11.5, 0.8, [(psub, 22, False, "#F1F3F4")])
    _txt(s, 0.9, 5.6, 11.5, 0.6,
         [(f"客户：{meta.get('client','')}   ｜   类型：{meta.get('video_type','')}   ｜   "
           f"时长：{meta.get('duration','')}   ｜   画幅：{meta.get('aspect_ratio','')}", 13, False, "#E8EAED")])
    if meta.get("logo_path") and os.path.exists(meta["logo_path"]):
        try:
            s.shapes.add_picture(meta["logo_path"], Inches(0.9), Inches(6.4), height=Inches(0.7))
        except Exception:
            pass


# ---- P2 目录 ---------------------------------------------------------------
def _toc(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "目录 · CONTENTS", accent)
    items = [
        ("01", "项目理解", "我们对需求的解读"),
        ("02", "核心策略", "为什么选这个创意方向"),
        ("03", "创意概念", "一句话创意与情绪关键词"),
        ("04", "故事架构", "叙事结构与情绪曲线"),
        ("05", "分镜精选", "关键镜头分镜图与说明"),
        ("06", "视觉风格", "Moodboard 与影调色板"),
        ("07", "美术置景", "主场景与服化道总览"),
        ("08", "演员角色", "角色设定与选角方向"),
        ("09", "制作排期", "筹备-拍摄-后期时间线"),
        ("10", "执行团队", "主创介绍与同类案例"),
    ]
    x0, y0, col_w, row_h = 0.9, 1.6, 5.9, 0.95
    for i, (num, t1, t2) in enumerate(items):
        col = i % 2
        row = i // 2
        l = x0 + col * col_w
        t = y0 + row * row_h
        _txt(s, l, t, 1.0, 0.8, [(num, 30, True, accent)])
        _txt(s, l + 1.0, t + 0.05, col_w - 1.0, 0.8,
             [[(t1, 16, True, "#202124")], [(t2, 11, False, "#80868B")]])
    _footer(s, project, 2)


# ---- P3 项目理解 -----------------------------------------------------------
def _p3_understanding(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "01 · 项目理解", accent)
    brd = project["stage1"].get("brd", {})
    items = [
        ("项目定位", brd.get("positioning", "（待补充）")),
        ("核心目标", brd.get("core_goal", "（待补充）")),
        ("目标受众", brd.get("audience", "（待补充）")),
        ("必传信息", brd.get("must_info", "（待补充）")),
        ("风格倾向", brd.get("style", "（待补充）")),
        ("约束条件", brd.get("constraints", "（待补充）")),
        ("隐性需求（推断）", brd.get("hidden", "（待补充）")),
    ]
    _bullet(s, 0.9, 1.6, 11.5, 5.0, items, accent, size=15, gap=10)
    _footer(s, project, 3)


# ---- P4 核心策略 -----------------------------------------------------------
def _p4_strategy(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "02 · 核心策略", accent)
    dirs = project["stage2"].get("directions", [])
    idx = min(project["stage2"].get("chosen", 0), max(len(dirs) - 1, 0))
    d = dirs[idx] if idx < len(dirs) else {}
    _txt(s, 0.9, 1.5, 11.5, 0.5,
         [("为什么选「" + (d.get("name", "该方向") or "该方向") + "」这个方向", 18, True, accent)])
    items = [
        ("一句话创意", d.get("logline", "（待补充）")),
        ("适配理由", d.get("rationale", "（待补充）")),
        ("目标受众契合", project["stage1"].get("brd", {}).get("audience", "（待补充）")),
        ("策略推导", "从“受众-渠道双约束“出发：本片投放在「" +
         (project["meta"].get("channel", "渠道") or "渠道") + "」，受众为「" +
         (project["stage1"].get("brd", {}).get("audience", "目标人群") or "目标人群") +
         "」，因此采用该叙事结构与视觉调性以最大化转化。"),
    ]
    _bullet(s, 0.9, 2.2, 11.5, 4.5, items, accent, size=15, gap=12)
    _footer(s, project, 4)


# ---- P5 创意概念 -----------------------------------------------------------
def _p5_concept(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "03 · 创意概念", accent)
    dirs = project["stage2"].get("directions", [])
    idx = min(project["stage2"].get("chosen", 0), max(len(dirs) - 1, 0))
    d = dirs[idx] if idx < len(dirs) else {}
    _txt(s, 0.9, 1.5, 11.5, 1.2,
         [[("一句话创意  ", 14, True, "#80868B"), (d.get("logline", "（待补充）"), 24, True, accent)]])
    items = [
        ("核心概念阐述", d.get("narrative", "（待补充）")),
        ("视觉调性", d.get("visual_tone", "（待补充）")),
        ("情绪曲线", d.get("emotion_curve", "（待补充）")),
    ]
    _bullet(s, 0.9, 3.0, 11.5, 3.5, items, accent, size=15, gap=12)
    _footer(s, project, 5)


# ---- P6 故事梗概 -----------------------------------------------------------
def _p6_story(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "04 · 故事架构", accent)
    segs = project["stage3"].get("segments", [])
    if not segs:
        _txt(s, 0.9, 1.6, 11.5, 4.5, [("（待补充：请在环节3完成创意脚本）", 16, False, "#80868B")])
    else:
        y = 1.6
        for i, seg in enumerate(segs, 1):
            _txt(s, 0.9, y, 1.4, 0.7, [(f"段落{i}", 16, True, accent)])
            _txt(s, 2.3, y, 10.1, 0.7,
                 [[(seg.get("time", "") + "  ", 12, True, "#80868B"),
                   (seg.get("visual", "（待补充）"), 14, False, "#202124")],
                  [("旁白：" + seg.get("voiceover", ""), 12, False, "#5F6368")]])
            y += 0.95
    _footer(s, project, 6)


# ---- P7-P9 分镜精选 --------------------------------------------------------
def _p7_9_storyboard(prs, project, accent):
    boards = project["stage5"].get("boards", [])
    shots = project["stage4"].get("shots", [])
    # 取前 9 个关键镜头，3 个一组
    chunks = [boards[i:i+3] for i in range(0, min(len(boards), 9), 3)] or [[]]
    page = 7
    for ci, chunk in enumerate(chunks):
        s = _blank(prs)
        _top_bar(s, f"05 · 分镜精选（{ci+1}/{len(chunks)}）", accent)
        n = len(chunk)
        card_w = 3.9
        gap = 0.35
        x0 = 0.9
        y0 = 1.5
        for j, bd in enumerate(chunk):
            l = x0 + j * (card_w + gap)
            global_idx = ci * 3 + j
            shot = shots[global_idx] if global_idx < len(shots) else {}
            # 图片或占位
            img_path = bd.get("image_path", "")
            if img_path and os.path.exists(img_path):
                try:
                    s.shapes.add_picture(img_path, Inches(l), Inches(y0), width=Inches(card_w))
                except Exception:
                    _rect(s, l, y0, card_w, 2.6, "#ECEFF1")
            else:
                _rect(s, l, y0, card_w, 2.6, "#ECEFF1")
                _txt(s, l, y0 + 1.0, card_w, 0.6,
                     [("分镜图占位", 13, False, "#9AA0A6")], align=PP_ALIGN.CENTER)
            # 说明
            cap = (f"镜{global_idx+1} · {shot.get('shot_size','')} · {shot.get('movement','')}\n"
                   f"{shot.get('desc','（待补充画面描述）')}")
            _txt(s, l, y0 + 2.75, card_w, 2.4,
                 [[(cap, 12, False, "#3C4043")]])
        _footer(s, project, page)
        page += 1
    # 若不足 3 页，补齐到至少 3 页占位
    while page <= 9:
        s = _blank(prs)
        _top_bar(s, f"05 · 分镜精选（{page-6}/3）", accent)
        _txt(s, 0.9, 3.0, 11.5, 1.0, [("（更多镜头见文字分镜表 / 动态分镜）", 14, False, "#80868B")])
        _footer(s, project, page)
        page += 1


# ---- P10 视觉风格 ----------------------------------------------------------
def _p10_style(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "06 · 视觉风格", accent)
    mb = project["stage7"].get("moodboard", {})
    items = [
        ("主色板", mb.get("color_palette", "（待补充）")),
        ("光影风格", mb.get("light_style", "（待补充）")),
        ("影调关键词", mb.get("tone_keywords", "（待补充）")),
        ("参考片推荐", mb.get("ref_films", "（待补充）")),
        ("摄影质感", mb.get("photo_texture", "（待补充）")),
    ]
    _bullet(s, 0.9, 1.6, 7.0, 5.0, items, accent, size=14, gap=10)
    # 右侧风格参考列表
    refs = project["stage7"].get("style_refs", [])
    _txt(s, 8.3, 1.6, 4.2, 0.5, [("风格参考片", 15, True, accent)])
    ref_items = [r for r in refs if r.strip()] or ["（待补充）"]
    _bullet(s, 8.3, 2.1, 4.2, 4.0, ref_items, "#80868B", size=12, gap=8)
    _footer(s, project, 10)


# ---- P11 美术置景 ----------------------------------------------------------
def _p11_art(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "07 · 美术置景", accent)
    art = project["stage7"].get("art_set", [])
    props = project["stage7"].get("props", [])
    y = 1.6
    _txt(s, 0.9, y, 11.5, 0.5, [("主场景概念", 16, True, accent)]); y += 0.55
    if art:
        for sc in art:
            _bullet(s, 0.9, y, 11.5, 0.9,
                    [(sc.get("name", "（场景）"),
                      f"  空间：{sc.get('space','')} ｜ 色调光线：{sc.get('tone','')} ｜ 预算：{sc.get('budget','')}")],
                    accent, size=13, gap=4)
            y += 0.95
    else:
        _txt(s, 0.9, y, 11.5, 0.5, [("（待补充）", 13, False, "#80868B")]); y += 0.6
    y += 0.3
    _txt(s, 0.9, y, 11.5, 0.5, [("核心道具", 16, True, accent)]); y += 0.55
    if props:
        for p in props:
            _bullet(s, 0.9, y, 11.5, 0.8,
                    [(p.get("name", "（道具）"),
                      f"  意义：{p.get('meaning','')} ｜ 筹备：{p.get('prep','')}")],
                    accent, size=13, gap=4)
            y += 0.85
    else:
        _txt(s, 0.9, y, 11.5, 0.5, [("（待补充）", 13, False, "#80868B")])
    _footer(s, project, 11)


# ---- P12 演员与角色 --------------------------------------------------------
def _p12_actors(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "08 · 演员与角色", accent)
    actors = project["stage7"].get("actors", [])
    if not actors:
        _txt(s, 0.9, 1.6, 11.5, 4.5, [("（待补充：请在环节7填写角色设定）", 16, False, "#80868B")])
    else:
        y = 1.5
        for a in actors:
            _txt(s, 0.9, y, 11.5, 0.5, [(a.get("name", "（角色）"), 16, True, accent)]); y += 0.5
            _bullet(s, 0.9, y, 11.5, 1.6,
                    [("外形", a.get("look", "（待补充）")),
                     ("演员参考", a.get("actor_ref", "（待补充）")),
                     ("服装", a.get("wardrobe", "（待补充）")),
                     ("妆发", a.get("makeup", "（待补充）"))],
                    "#80868B", size=12, gap=3)
            y += 1.9
    _footer(s, project, 12)


# ---- P13 制作排期 ----------------------------------------------------------
def _p13_schedule(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "09 · 制作排期", accent)
    phases = [
        ("筹备期", "需求确认 / 创意定稿 / 分镜与动态分镜 / 选角置景 / 服化道筹备", "第 1-2 周"),
        ("拍摄期", "主场景拍摄 / 补拍 / 素材整理", "第 3 周"),
        ("后期期", "剪辑 / 调色 / 配音配乐 / 特效 / 成片交付", "第 4-5 周"),
    ]
    x = 0.9
    col_w = 3.7
    for i, (name, desc, time) in enumerate(phases):
        l = x + i * (col_w + 0.4)
        _rect(s, l, 1.7, col_w, 0.7, accent)
        _txt(s, l, 1.7, col_w, 0.7, [(name, 18, True, "#FFFFFF")],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, l, 2.6, col_w, 2.6, [[(desc, 13, False, "#3C4043")]])
        _txt(s, l, 5.3, col_w, 0.5, [(time, 13, True, accent)], align=PP_ALIGN.CENTER)
    # 用户可编辑排期备注
    extra = project["stage8"].get("schedule_text", "")
    if extra:
        _txt(s, 0.9, 6.0, 11.5, 0.8, [[("备注：" + extra, 12, False, "#5F6368")]])
    _footer(s, project, 13)


# ---- P14 团队与优势 --------------------------------------------------------
def _p14_team(prs, project, accent):
    s = _blank(prs)
    _top_bar(s, "10 · 执行团队", accent)
    team = project["stage8"].get("team_text", "")
    default_team = "（请在 PPT 整合页填写主创介绍与同类案例）"
    _txt(s, 0.9, 1.6, 11.5, 4.5,
         [[(team or default_team, 14, False, "#3C4043")]])
    _footer(s, project, 14)


# ---- P15 结尾 --------------------------------------------------------------
def _p15_end(prs, project, accent):
    s = _blank(prs)
    _rect(s, 0, 0, 13.333, 7.5, accent)
    dirs = project["stage2"].get("directions", [])
    idx = min(project["stage2"].get("chosen", 0), max(len(dirs) - 1, 0))
    slogan = dirs[idx].get("logline", "让创意，成为生意。") if idx < len(dirs) else "让创意，成为生意。"
    contact = project["stage8"].get("contact", "")
    _txt(s, 0.9, 2.6, 11.5, 1.5, [(slogan, 34, True, "#FFFFFF")])
    if contact:
        _txt(s, 0.9, 4.3, 11.5, 0.8, [(contact, 16, False, "#F1F3F4")])
    _txt(s, 0.9, 6.4, 11.5, 0.5,
         [(project["meta"].get("project_name", "商业视频创意方案"), 12, False, "#E8EAED")])
